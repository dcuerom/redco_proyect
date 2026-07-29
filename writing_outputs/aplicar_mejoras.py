# -*- coding: utf-8 -*-
"""Aplica los 'Puntos a mejorar' de la tarea Notion sobre REDCO_Presentacion_Historico - Respaldo.pptx"""
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

SRC = "/sessions/determined-fervent-noether/mnt/redco_proyect/REDCO_Presentacion_Historico  -  Respaldo.pptx"
OUT = "/sessions/determined-fervent-noether/mnt/redco_proyect/REDCO_Presentacion_Historico_rev1.pptx"

NAVY   = RGBColor(0x1E, 0x27, 0x61)
BLUE   = RGBColor(0x2C, 0x4A, 0x7C)
GRAY   = RGBColor(0x6B, 0x72, 0x80)
DARK   = RGBColor(0x23, 0x27, 0x33)
LINE   = RGBColor(0xE2, 0xE5, 0xEA)
CARD   = RGBColor(0xF4, 0xF6, 0xF9)
BAND   = RGBColor(0xEE, 0xF2, 0xF7)
ORANGE = RGBColor(0xC8, 0x86, 0x2D)
GREEN  = RGBColor(0x2E, 0x7D, 0x5B)
RED    = RGBColor(0xB4, 0x44, 0x3C)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Calibri"


# ---------- helpers ----------
def set_text(shape, text, size=None, bold=None, color=None, align=None, italic=None):
    """Reescribe el texto de una caja conservando el formato del primer run."""
    tf = shape.text_frame
    p = tf.paragraphs[0]
    runs = p.runs
    if not runs:
        r = p.add_run()
    else:
        r = runs[0]
        for extra in runs[1:]:
            extra._r.getparent().remove(extra._r)
    r.text = text
    for extra_p in tf.paragraphs[1:]:
        extra_p._p.getparent().remove(extra_p._p)
    f = r.font
    f.name = FONT
    if size is not None:
        f.size = Pt(size)
    if bold is not None:
        f.bold = bold
    if italic is not None:
        f.italic = italic
    if color is not None:
        f.color.rgb = color
    if align is not None:
        p.alignment = align
    return shape


def txbox(slide, l, t, w, h, text, size=9, bold=False, color=DARK,
          align=PP_ALIGN.LEFT, italic=False, anchor=MSO_ANCHOR.TOP, spacing=None):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if spacing:
        p.line_spacing = spacing
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def add_run(tb, text, size=9, bold=False, color=DARK, italic=False):
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return r


def bullets(slide, l, t, w, h, items, size=8, color=DARK):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for k, item in enumerate(items):
        p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        p.line_spacing = 0.95
        p.space_after = Pt(3)
        rb = p.add_run()
        rb.text = "•  "
        rb.font.name, rb.font.size, rb.font.bold = FONT, Pt(size), True
        rb.font.color.rgb = color
        r = p.add_run()
        r.text = item
        r.font.name, r.font.size = FONT, Pt(size)
        r.font.color.rgb = color
    return tb


def rect(slide, l, t, w, h, fill, shape=MSO_SHAPE.RECTANGLE, line=None, adj=None):
    sp = slide.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(0.75)
    sp.shadow.inherit = False
    if adj is not None:
        try:
            sp.adjustments[0] = adj
        except Exception:
            pass
    tf = sp.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return sp


def send_to_back(slide, shape):
    tree = slide.shapes._spTree
    el = shape._element
    tree.remove(el)
    tree.insert(2, el)


def by_id(slide, sid):
    for sh in slide.shapes:
        if sh.shape_id == sid:
            return sh
    raise KeyError(sid)


def by_name(slide, name):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    raise KeyError(name)


prs = Presentation(SRC)
S = prs.slides


# =====================================================================
# SLIDE 8 — leyenda de fuente: quitar la indicación de "Colores..."
# =====================================================================
s8 = S[7]
set_text(by_name(s8, "TextBox 13"),
         "Fuente: 202601_Proyeccion_Ing+Modulos 1_Edu.xlsx (Gantt Proyectos + Listas y Parametros) "
         "y 00_Ciclo Mensual EdP 2023_v0.xlsx (Ingresado).")

# =====================================================================
# SLIDE 9 — leyenda de fuente: quitar la indicación de "Colores..."
# =====================================================================
s9 = S[8]
set_text(by_name(s9, "TextBox 22"),
         "Fuente: 202601_Proyeccion_Ing+Modulos 1_Edu.xlsx (Gantt + Listas) y "
         "00_Ciclo Mensual EdP 2023_v0.xlsx (Ingresado). 32 proyectos con código mapeado; "
         "rent. = ingreso cobrado − costo modelado (tarifas estándar).")


# =====================================================================
# SLIDE 13 — cambio % entre años + nota de diversificación
# =====================================================================
s13 = S[12]

# corrección del título ("Producción por  por país")
set_text(by_id(s13, 5), "Producción por país — composición anual")

# --- deltas interanuales dentro de la dona ---
# totales EDP Emitido USD (de los charts): 2023 7,071,137 · 2024 8,418,845 · 2025 9,476,037
deltas = [
    (None, None),
    (5.67, "▲ +19,1%  vs 2023"),
    (9.79, "▲ +12,6%  vs 2024"),
]
for left, label in deltas:
    if label is None:
        continue
    tb = txbox(s13, left - 0.30, 4.16, 2.60, 0.20, label,
               size=9, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# --- nota de diversificación ---
band = rect(s13, 0.60, 6.34, 12.14, 0.64, CARD, MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.10)
accent = rect(s13, 0.60, 6.34, 0.06, 0.64, ORANGE)
txbox(s13, 0.80, 6.44, 2.60, 0.16, "NOTA · DIVERSIFICACIÓN",
      size=8, bold=True, color=ORANGE)
nb = txbox(s13, 0.80, 6.63, 11.80, 0.30,
           "El crecimiento del período (+34,0 % ’23→’25) no viene de Chile: Chile cae de $4,3M a $3,3M "
           "(−24 %) y su peso baja de 61 % a 34 %. ",
           size=9, color=DARK)
add_run(nb, "Perú (+52 %), Brasil (×2,5) y USA (×16) absorben todo el crecimiento; "
            "los países efectivos (1/HHI) pasan de 2,3 a 3,9.", size=9, bold=True, color=NAVY)


# =====================================================================
# SLIDE 14 — separación visual entre PRODUCCIÓN y DOTACIÓN
# =====================================================================
s14 = S[13]

# 1) comprimir el bloque de dotación para dejar aire bajo el rótulo de sección
BASE = 6.28          # línea base de las barras de dotación
K = 0.86             # factor de compresión
for sh in list(s14.shapes):
    t = Emu(sh.top).inches
    h = Emu(sh.height).inches
    w = Emu(sh.width).inches
    if 4.55 <= t <= 6.30 and w <= 0.90:      # barras (w=0.62) y etiquetas (w<=0.86)
        new_t = BASE - (BASE - t) * K
        sh.top = Inches(new_t)
        if abs(w - 0.62) < 0.01:             # sólo las barras cambian de alto
            sh.height = Inches(max(h * K, 0.02))

# 2) banda de fondo para la sección DOTACIÓN
band14 = rect(s14, 0.30, 4.24, 12.73, 2.34, BAND, MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.03)
send_to_back(s14, band14)

# 3) rótulos de sección con el mismo lenguaje visual que los chips de producción
#    (chip azul marino = dotación · chip de color = producción)
dot_chips = {
    44:  (0.73, "Dotación  55 → 76"),     # Chile
    74:  (3.31, "Dotación  21 → 30"),     # Perú
    103: (5.89, "Dotación  4 → 3"),       # Brasil
    126: (8.47, "Dotación  1 → 3"),       # USA
    148: (11.05, "Dotación  0 → 1"),      # Rusia
}
for sid, (x, label) in dot_chips.items():
    chip = rect(s14, x, 4.36, 1.70, 0.30, NAVY, MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.42)
    lbl = by_id(s14, sid)
    lbl.left, lbl.top, lbl.width, lbl.height = Inches(x), Inches(4.44), Inches(1.70), Inches(0.15)
    set_text(lbl, label, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # el chip debe quedar detrás del texto
    tree = s14.shapes._spTree
    tree.remove(chip._element)
    tree.insert(list(tree).index(lbl._element), chip._element)

# 4) leyenda de categorías: bajarla para que no toque el borde de la banda
for sid in (154, 155, 156, 157, 158, 159, 160, 161, 162, 163):
    sh = by_id(s14, sid)
    sh.top = Inches(Emu(sh.top).inches + 0.07)
txbox(s14, 0.50, 6.72, 0.90, 0.24, "Categorías:", size=9, bold=True, color=GRAY)


# =====================================================================
# SLIDE 11 — rediseño: ticket promedio con variación interanual,
#            embudo emitido→adjudicado y bloque de lecturas
# =====================================================================
s11 = S[10]

KEEP = {"TextBox 1", "Oval 2", "Oval 3", "TextBox 4", "TextBox 5", "TextBox 56"}
for sh in list(s11.shapes):
    if sh.name not in KEEP:
        sh._element.getparent().remove(sh._element)

# --- encabezado ---
title = by_name(s11, "TextBox 4")
title.top, title.height = Inches(0.83), Inches(0.44)
set_text(title, "Ventas: ticket promedio y conversión de propuestas")

sub = by_name(s11, "TextBox 5")
sub.left, sub.top, sub.width, sub.height = Inches(0.98), Inches(1.42), Inches(11.80), Inches(0.23)
set_text(sub, "Se ganan las propuestas chicas: por monto se adjudica ~1 de cada 4–5 USD emitidos, "
              "y el ticket adjudicado se contrae año a año.")

# --- leyenda de fuente: 2024 es año completo, no sólo Q4 ---
set_text(by_name(s11, "TextBox 56"),
         "Fuente: 01_KPI Gestión REDCO_2026.xlsx (KPI´s de negocio 2024/2025/2026). "
         "Ticket = monto ÷ N° de propuestas.  2024 = año completo (ene–dic);  2026 * = enero–junio.")
by_name(s11, "TextBox 56").top = Inches(7.04)

COLS = [0.42, 4.63, 8.84]
CW = 4.06

# separadores verticales entre columnas-año
for x in (4.555, 8.765):
    rect(s11, x, 1.86, 0.012, 3.95, LINE)

YEARS = [
    # etiqueta, N emitidas, monto emitido M, ticket emitido k, N adj, monto adj M,
    # ticket adj k, conv monto, conv N°, Δ ticket emitido, Δ ticket adjudicado
    ("2024",   119, "34.2", "$289k", 53, "12.2", "$230k", 0.36, "53/119 (45%)", None,   None),
    ("2025",   164, "43.7", "$266k", 78, "10.3", "$132k", 0.24, "78/164 (48%)", "−8%",  "−43%"),
    ("2026 *",  86, "19.1", "$222k", 34,  "4.1", "$119k", 0.21, "34/86 (40%)",  "−17%", "−10%"),
]
CONV_LBL = ["36%", "24%", "21%"]
# ventana comparable ene–jun (fuente: KPI´s de negocio 2025 / 2026, suma mensual ene–jun)
COMPARABLE = [
    "",
    "ventana comparable ene–jun: 41/76 (54%) · 21,7% del monto",
    "* 2026 es, por definición, la ventana ene–jun",
]

for i, (yr, n_em, m_em, t_em, n_ad, m_ad, t_ad, conv, byn, d_em, d_ad) in enumerate(YEARS):
    x = COLS[i]

    # año
    txbox(s11, x, 1.86, CW, 0.32, yr, size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    rect(s11, x, 2.26, CW, 0.02, LINE)

    # tarjetas de ticket
    for j, (lbl, val, det, delta) in enumerate([
        ("TICKET EMITIDO", t_em, f"{n_em} prop. · ${m_em}M", d_em),
        ("TICKET ADJUDICADO", t_ad, f"{n_ad} adj. · ${m_ad}M", d_ad),
    ]):
        cx = x + j * 2.10
        rect(s11, cx, 2.30, 1.96, 1.14, CARD, MSO_SHAPE.ROUNDED_RECTANGLE, line=LINE, adj=0.06)
        txbox(s11, cx + 0.14, 2.40, 1.68, 0.16, lbl, size=8.5, bold=True, color=GRAY)
        txbox(s11, cx + 0.14, 2.58, 1.15, 0.50, val, size=26, bold=True, color=BLUE)
        if delta:
            chip = rect(s11, cx + 1.22, 2.72, 0.60, 0.22, RED, MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.45)
            tf = chip.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = delta
            r.font.name, r.font.size, r.font.bold = FONT, Pt(9), True
            r.font.color.rgb = WHITE
        else:
            txbox(s11, cx + 1.22, 2.76, 0.60, 0.18, "base", size=8.5, italic=True,
                  color=GRAY, align=PP_ALIGN.CENTER)
        txbox(s11, cx + 0.14, 3.16, 1.68, 0.16, det, size=8.5, color=DARK)

    txbox(s11, x, 3.52, CW, 0.16, "Del monto emitido, cuánto se adjudica",
          size=8.5, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

    # embudo
    bar = rect(s11, x, 3.70, CW, 0.42, BLUE)
    tfb = bar.text_frame
    tfb.vertical_anchor = MSO_ANCHOR.MIDDLE
    pb = tfb.paragraphs[0]
    pb.alignment = PP_ALIGN.CENTER
    rb = pb.add_run()
    rb.text = f"Emitido  ${m_em}M"
    rb.font.name, rb.font.size, rb.font.bold = FONT, Pt(15), True
    rb.font.color.rgb = WHITE

    rect(s11, x + CW / 2 - 0.13, 4.18, 0.26, 0.24, GREEN, MSO_SHAPE.DOWN_ARROW)
    txbox(s11, x, 4.45, CW, 0.15, "Adjudicado", size=8, bold=True, color=NAVY,
          align=PP_ALIGN.CENTER)

    aw = CW * conv
    abar = rect(s11, x + (CW - aw) / 2, 4.64, aw, 0.42, NAVY)
    tfa = abar.text_frame
    tfa.vertical_anchor = MSO_ANCHOR.MIDDLE
    pa = tfa.paragraphs[0]
    pa.alignment = PP_ALIGN.CENTER
    ra = pa.add_run()
    ra.text = f"${m_ad}M"
    ra.font.name, ra.font.size, ra.font.bold = FONT, Pt(13), True
    ra.font.color.rgb = WHITE

    conv_tb = txbox(s11, x, 5.12, CW, 0.30, CONV_LBL[i], size=20, bold=True,
                    color=NAVY, align=PP_ALIGN.CENTER)
    add_run(conv_tb, " del monto emitido", size=9.5, color=GRAY)
    txbox(s11, x, 5.43, CW, 0.16, f"por N°: {byn}", size=9, color=GRAY, align=PP_ALIGN.CENTER)
    if COMPARABLE[i]:
        txbox(s11, x, 5.61, CW, 0.15, COMPARABLE[i], size=8, italic=True,
              color=ORANGE if i == 1 else GRAY, align=PP_ALIGN.CENTER)

# --- bloque de lecturas / inferencias ---
LY, LH = 5.86, 1.08
LECTURAS = [
    (0.42, NAVY, "LECTURA POR CANTIDAD DE PROPUESTAS", [
        "La conversión por número de propuestas cayó en 2026 vs. la ventana comparable de 2025 "
        "(54 % → 40 %): en 2026 se emiten más propuestas pero se adjudica una fracción menor.",
        "La conversión por valor es mucho más baja que por cantidad (~21–24 % vs ~48 %): REDCO gana "
        "muchas propuestas chicas pero convierte poco del monto — las propuestas de mayor ticket caen más.",
    ]),
    (6.78, ORANGE, "LECTURA POR TRANSFORMACIÓN DE MONTO", [
        "REDCO adjudica propuestas de la mitad del tamaño que emite.",
        "Baja sostenida del ticket promedio emitido.",
        "Mejora leve del ratio Adj/Emit porque en 2026 bajó el tamaño promedio emitido, "
        "acercándolo al adjudicado.",
    ]),
]
for lx, acc, head, items in LECTURAS:
    rect(s11, lx, LY, 6.12, LH, CARD, MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.07)
    rect(s11, lx, LY, 0.05, LH, acc)
    txbox(s11, lx + 0.18, LY + 0.11, 5.70, 0.15, head, size=8, bold=True, color=acc)
    bullets(s11, lx + 0.18, LY + 0.33, 5.76, LH - 0.44, items, size=8, color=DARK)

prs.save(OUT)
print("OK ->", OUT)
