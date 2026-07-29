#!/usr/bin/env python3
"""Construye REDCO_Presentacion_Estrategica_2026.pptx.

Rehace los capítulos 2, 3 y 4 de `240626 Proyecto Gestión Estratégica REDCO(2).pptx`
con la data de `pilar_a` y las cifras de mercado actualizadas, bajo el
`Estandar_Presentaciones_PPT_REDCO(1).md`.

Uso:  .venv/bin/python pilar_a/scripts/build_presentacion_estrategica.py
"""

import math
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

sys.path.insert(0, str(Path(__file__).resolve().parent))

from redco_datos import cargar  # noqa: E402
from redco_deck import (  # noqa: E402
    AMBER, BG, GREEN, GREY, ICE, INK, LINE, MID, NAVY, RED, SERIES_RAMP, SKY,
    SLATE, STEEL, WHITE, BRAND_FONT, FONT, MARGIN, SLIDE_H, SLIDE_W, FOOTER_Y,
    add_chart, card, cover, divider, header, insight, line_h, oval, rect,
    shape_text, source_note, txt, txt_inline,
)

ROOT = Path(__file__).resolve().parents[2]
OUT = ROOT / "REDCO_Presentacion_Estrategica_2026.pptx"

FUENTE_PILAR_A = ("Fuente: pilar_a · REDCO_Rentabilidad_Ventas_FlujoCaja_historico.xlsx "
                  "(ledger CicloEdP_2023_v1, cierre 2025 y KPI 2026).")


# ============================================================ infraestructura

def nueva_presentacion():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    _declarar_fuente_marca(prs)
    return prs


def _declarar_fuente_marca(prs):
    """Estándar §4.1: declarar Lexend Deca en el tema, con Calibri como fallback
    explícito en los runs (Lexend Deca no está instalada en el entorno)."""
    from lxml import etree
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    try:
        theme = prs.slide_masters[0].part.part_related_by(
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme")
        root = etree.fromstring(theme.blob)
        cambios = 0
        for tag in ("majorFont", "minorFont"):
            for f in root.iter(f"{ns}{tag}"):
                for latin in f.findall(f"{ns}latin"):
                    latin.set("typeface", BRAND_FONT)
                    cambios += 1
        theme._blob = etree.tostring(root, xml_declaration=True,
                                     encoding="UTF-8", standalone=True)
        if not cambios:
            raise ValueError("no se encontró fontScheme/latin")
    except Exception as exc:  # el deck sigue siendo válido sin esto
        print(f"  aviso: no se pudo declarar {BRAND_FONT} en el tema ({exc})")


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def fmt_musd(v, dec=1):
    return f"{v / 1e6:,.{dec}f}".replace(",", "@").replace(".", ",").replace("@", ".")


def pct(v, dec=0):
    return f"{v * 100:,.{dec}f}".replace(".", ",") + "%"


# ================================================================== CAPÍTULO 2

def s01_portada(prs):
    cover(blank(prs),
          "Gestión Estratégica REDCO",
          "Actualización 2026 · Resultados, metodología de decisión y metas",
          "Capítulos 2 – 4  ·  Horizonte REDCO 2028",
          "Documento interno · Data: pilar_a (ciclo EdP 2023–2026) · Mercado: ETC, "
          "Wood Mackenzie e IEA")


def s02_agenda(prs):
    s = blank(prs)
    header(s, "Contenido",
           "Tres bloques: dónde estamos, cómo decidimos dónde crecer, y qué comprometemos.")
    temas = [
        ("Primeros resultados obtenidos",
         "La evidencia 2023–2026 de rentabilidad, venta y cartera, y lo que el "
         "mercado minero proyecta a 2035."),
        ("Metodología para la toma de decisiones",
         "El marco de atractivo de mercado y habilidad de ganar, y la disciplina "
         "de distancia al CORE."),
        ("Propósito, misión y habilitadores",
         "El propósito y la misión vigentes, los motores de crecimiento y las "
         "ocho metas al 2028."),
    ]
    y = 2.05
    for i, (titulo, bajada) in enumerate(temas, 1):
        h = 1.30
        rect(s, MARGIN, y, SLIDE_W - 2 * MARGIN, h, fill=BG, line=LINE)
        rect(s, MARGIN, y, 0.05, h, fill=STEEL)
        o = oval(s, MARGIN + 0.36, y + 0.42, 0.46, fill=STEEL)
        shape_text(o, str(i), size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, pad=0)
        txt(s, MARGIN + 1.06, y + 0.30, 9.9, 0.34, titulo, size=16, bold=True, color=INK)
        txt(s, MARGIN + 1.06, y + 0.70, 10.6, 0.44, bajada, size=11.5, color=GREY)
        y += h + 0.24
    source_note(s, "Los tres bloques corresponden 1:1 con los capítulos 2, 3 y 4 del cuerpo "
                   "del documento: ningún tema queda fuera y ningún ítem está de más.")


def s03_divisor_cap2(prs):
    divider(blank(prs), "02", "Primeros resultados obtenidos",
            "Qué muestra la data real de REDCO 2023–2026 y qué proyecta el mercado minero "
            "hacia 2035.")


def s04_contexto(prs, d):
    s = blank(prs)
    header(s, "REDCO creció y se diversificó; falta sostener la rentabilidad",
           "El relato 2024 sigue siendo válido en dirección, pero la data 2025 cambia el énfasis.")

    txt(s, MARGIN, 1.72, 7.15, 2.60,
        [("REDCO, consultora boutique especializada en Consultoría, Estudios e Ingeniería "
          "minera, lleva más de 20 años en el mercado con presencia en Chile, Perú, Brasil, "
          "USA y operaciones intermitentes en Rusia.", {"size": 12.5, "spacing": 9}),
         ("Lo que en 2024 se describía como “un nuevo mercado como lo es EEUU” hoy es una "
          "operación consolidada: USA aporta 1,65 MUSD y es el cuarto país por venta. "
          "El equipo pasó de 81 a 113 personas en un año.", {"size": 12.5, "spacing": 9}),
         ("El mercado minero —y en particular los minerales de la transición energética— "
          "sigue siendo un escenario promisorio. La restricción ya no es la demanda: "
          "es la capacidad de ejecutar con rentabilidad.", {"size": 12.5})],
        color=INK)

    tiles = [
        ("VENTA 2025", fmt_musd(d["venta_total"][2]) + " MUSD", "+34% vs 2023", GREEN),
        ("PAÍSES EFECTIVOS", str(d["paises_efectivos"][2025]).replace(".", ","),
         "desde 2,3 en 2023", GREEN),
        ("MARGEN OPERACIONAL", pct(d["consolidado4"]["2025"]["margen_pct"], 1),
         "desde 26,5% en 2024", RED),
        ("DOTACIÓN", f"{d['dotacion_cat']['TOTAL']['2026']}", "desde 81 en 2025", AMBER),
    ]
    x = 7.85
    for i, (label, valor, delta, color) in enumerate(tiles):
        y = 1.72 + i * 0.86
        rect(s, x, y, SLIDE_W - MARGIN - x, 0.74, fill=BG, line=LINE)
        rect(s, x, y, 0.05, 0.74, fill=color)
        txt(s, x + 0.18, y + 0.09, 3.2, 0.20, label, size=8.5, bold=True, color=GREY)
        txt(s, x + 0.18, y + 0.30, 2.3, 0.36, valor, size=20, bold=True, color=INK)
        txt(s, x + 2.45, y + 0.36, 2.2, 0.26, delta, size=9.5, color=color,
            align=PP_ALIGN.RIGHT)

    insight(s, MARGIN, 4.86, 7.15, 1.32, "LO QUE CAMBIÓ RESPECTO DE 2024", None,
            "El crecimiento se consiguió: la venta subió 34% en tres años y la dependencia "
            "de Chile bajó de 61% a 34%. Lo que no acompañó fue el resultado: la utilidad "
            "operacional se redujo a la mitad. Por eso este capítulo parte por la "
            "rentabilidad y no por la venta.", accent=NAVY)
    source_note(s, FUENTE_PILAR_A + "  Dotación: Dotación_Comparativo_Estandarizado_2025-2026.xlsx. "
                   "Margen operacional sobre perímetro de 4 países (Chile+Perú+Brasil+USA).")


def s05_crecimiento(prs, d):
    s = blank(prs)
    header(s, "La venta creció, pero la utilidad operacional se redujo a la mitad",
           "Dos paneles: el resultado del negocio y el tamaño de la organización que lo produce.")

    txt(s, MARGIN, 1.58, 6.0, 0.26, "Ingreso, margen y utilidad operacional",
        size=13, bold=True, color=INK)
    line_h(s, MARGIN, 1.90, 6.0, color=LINE)

    c4 = d["consolidado4"]
    ch = add_chart(s, XL_CHART_TYPE.COLUMN_CLUSTERED, MARGIN, 2.02, 6.0, 3.02,
                   ["2024", "2025"],
                   [("Ingreso Op. [MUSD]", [c4["2024"]["ingreso"] / 1e6, c4["2025"]["ingreso"] / 1e6]),
                    ("Margen Bruto [MUSD]", [c4["2024"]["margen_usd"] / 1e6, c4["2025"]["margen_usd"] / 1e6])],
                   colors=[NAVY, STEEL], legend=True, gap=90)
    ch.plots[0].has_data_labels = True
    dl = ch.plots[0].data_labels
    dl.number_format = '0.00'
    dl.number_format_is_linked = False
    dl.font.size = Pt(9)
    dl.font.color.rgb = INK

    for i, a in enumerate(("2024", "2025")):
        col = RED if c4[a]["margen_pct"] < 0.20 else GREEN
        bx = MARGIN + 1.30 + i * 2.75
        rect(s, bx, 5.14, 1.55, 0.44, fill=WHITE, line=col)
        txt(s, bx, 5.20, 1.55, 0.32, f"Utilidad  {pct(c4[a]['margen_pct'], 1)}",
            size=11.5, bold=True, color=col, align=PP_ALIGN.CENTER)

    txt(s, 6.98, 1.58, 5.85, 0.26, "Dotación y productividad", size=13, bold=True, color=INK)
    line_h(s, 6.98, 1.90, 5.85, color=LINE)

    cats = d["dotacion_cats"]
    add_chart(s, XL_CHART_TYPE.COLUMN_STACKED, 6.98, 2.02, 5.85, 3.02,
              ["2025", "2026"],
              [(c, [d["dotacion_cat"][c]["2025"], d["dotacion_cat"][c]["2026"]]) for c in cats],
              colors=[SLATE, SKY, MID, GREY, NAVY], legend=True, gap=110, overlap=100)

    ixp = d["ingreso_x_persona"]
    for i, (anio, val) in enumerate(sorted(ixp.items())):
        col = GREEN if anio == 2025 else RED
        bx = 8.30 + i * 2.75
        rect(s, bx, 5.14, 1.90, 0.44, fill=WHITE, line=col)
        txt(s, bx, 5.20, 1.90, 0.32, f"{anio} · ${val / 1000:,.0f}k / persona".replace(",", "."),
            size=11, bold=True, color=col, align=PP_ALIGN.CENTER)

    insight(s, MARGIN, 5.70, SLIDE_W - 2 * MARGIN, 1.24, "LECTURA",
            f"{pct(c4['2024']['margen_pct'], 1)} → {pct(c4['2025']['margen_pct'], 1)}",
            "La organización creció 40% en dotación (81 → 113) mientras la utilidad operacional "
            "caía a la mitad. A venta constante, el ingreso por persona baja de $113k a $81k: "
            "la nueva capacidad todavía no está pagada por venta nueva.", accent=RED)
    source_note(s, FUENTE_PILAR_A + "  Perímetro: 4 países (Chile+Perú+Brasil+USA), igual que el "
                   "análisis histórico. Utilidad = Margen ÷ Ingreso. Ingreso por persona = ingreso "
                   "operacional 2025 ÷ dotación del año. Nota: la serie 'dotación senior' del deck "
                   "2024 no es reproducible — el comparativo 2026 solo trae categoría, no cargo.")


def s06_ventas(prs, d):
    s = blank(prs)
    header(s, "La receta cambió: de 3–4 clientes clave a una cartera diversificada",
           "Venta por país y concentración de clientes a nivel grupo, 2023–2026.")

    paises = ["Chile", "Perú", "Brasil", "USA", "Rusia", "BVI", "Canadá"]
    cats = ["2023", "2024", "2025", "2026 *"]
    add_chart(s, XL_CHART_TYPE.COLUMN_STACKED, MARGIN, 1.62, 5.55, 2.92, cats,
              [(p, [v / 1e6 for v in d["venta_pais"][p]]) for p in paises],
              colors=SERIES_RAMP[:7], legend=True,
              gap=80, overlap=100)
    for i, tot in enumerate(d["venta_total"]):
        txt(s, MARGIN + 0.30 + i * 1.28, 4.46, 1.20, 0.22,
            f"{fmt_musd(tot)} M", size=10, bold=True, color=INK, align=PP_ALIGN.CENTER)

    txt(s, 6.42, 1.62, 6.4, 0.24, "Composición de la venta por cliente (grupo)",
        size=12, bold=True, color=INK)
    for i, anio in enumerate((2023, 2024, 2025)):
        tc = d["top_clientes"][anio]
        x = 6.42 + i * 2.16
        add_chart(s, XL_CHART_TYPE.DOUGHNUT, x - 0.12, 1.94, 2.30, 1.72,
                  [t[0] for t in tc["top4"]] + [tc["resto"][0]],
                  [("", [t[1] for t in tc["top4"]] + [tc["resto"][1]])],
                  point_colors=[NAVY, STEEL, MID, SKY, LINE], legend=False)
        txt(s, x + 0.14, 3.66, 1.80, 0.22, str(anio), size=11, bold=True,
            color=INK, align=PP_ALIGN.CENTER)
        txt(s, x + 0.14, 3.88, 1.80, 0.20,
            f"top-1  {tc['top1_pct']}%   ·   top-3  {tc['top3_pct']}%",
            size=8.5, color=GREY, align=PP_ALIGN.CENTER)
        txt(s, x + 0.14, 4.08, 1.80, 0.20,
            f"{tc['n_clientes']} clientes  ·  efect. {str(tc['clientes_efectivos']).replace('.', ',')}",
            size=8.5, color=GREY, align=PP_ALIGN.CENTER)

    labels = [(t[0], None) for t in d["top_clientes"][2025]["top4"]]
    txt(s, 6.42, 4.38, 6.4, 0.22,
        "Top 2025:  " + "   ·   ".join(f"{c} {p}%" for c, _, p in d["top_clientes"][2025]["top4"]),
        size=9, color=GREY)

    insight(s, MARGIN, 4.86, SLIDE_W - 2 * MARGIN, 2.04, "LO QUE ESTO SIGNIFICA",
            "32% → 12%",
            "En 2023 un solo cliente (CMP) explicaba el 32% de la venta del grupo y el top-3 el 56%. "
            "En 2025 el mayor cliente (BHP) pesa 12% y el top-3 solo 30%: los clientes efectivos "
            "pasan de 6,6 a 18,5 y los países efectivos de 2,3 a 3,9. La concentración que el deck "
            "2024 describía como “la receta” ya no es el motor — el crecimiento vino de diversificar. "
            "Riesgo a vigilar: en 2026 parcial, Rusia reaparece con 0,75 MUSD emitidos que no "
            "convierten a caja (efecto Seligdar).", accent=NAVY)
    source_note(s, FUENTE_PILAR_A + "  VENTA = 'Emitido USD' del ledger, atribuido por año de "
                   "proyecto (criterio que reproduce la hoja 02_Ventas). Clientes consolidados "
                   "(Cliente_Cons). Clientes/países efectivos = 1/HHI.  * 2026 = enero–junio.")


def s07_rentabilidad_pais(prs, d):
    s = blank(prs)
    header(s, "La rentabilidad del grupo vive fuera del país-sede",
           "Un panel por país en MUSD, todos a la misma escala para que sean comparables.")

    r = d["rentabilidad"]
    for i, p in enumerate(["Chile", "Perú", "Brasil", "USA"]):
        x = MARGIN + i * 3.15
        m24, m25 = r[p].get("margen_2024", 0), r[p].get("margen_2025", 0)
        col = RED if m25 < 0 else (GREEN if m25 >= 0.25 else AMBER)
        rect(s, x, 1.62, 2.95, 0.40, fill=NAVY)
        txt(s, x, 1.70, 2.95, 0.26, p, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_chart(s, XL_CHART_TYPE.COLUMN_CLUSTERED, x, 2.12, 2.95, 2.34, ["2024", "2025"],
                  [("Ingreso", [r[p].get("ing_2024", 0) / 1e6, r[p].get("ing_2025", 0) / 1e6]),
                   ("Margen", [r[p].get("benef_2024", 0) / 1e6, r[p].get("benef_2025", 0) / 1e6])],
                  colors=[NAVY, STEEL], legend=(i == 0), gap=80, font_size=8,
                  y_min=-2, y_max=4)
        rect(s, x, 4.56, 2.95, 0.60, fill=BG, line=LINE)
        txt(s, x + 0.10, 4.63, 2.75, 0.22, "UTILIDAD SOBRE INGRESO", size=7.5, bold=True,
            color=GREY, align=PP_ALIGN.CENTER)
        txt(s, x + 0.10, 4.84, 2.75, 0.26,
            f"{pct(m24, 0)}  →  {pct(m25, 0)}", size=13, bold=True, color=col,
            align=PP_ALIGN.CENTER)

    insight(s, MARGIN, 5.36, SLIDE_W - 2 * MARGIN, 1.52, "INSIGHT CLAVE",
            "−39%",
            "Chile concentra el mayor ingreso del grupo (3,59 MUSD) pero opera en pérdida, y la "
            "pérdida se profundiza: de −12% en 2024 a −39% en 2025. Absorbe la nómina de la "
            "operación regional — los sueldos son el 75% de su costo total. Perú (28%), USA (63%) "
            "y Brasil (72%) sostienen la rentabilidad del grupo. Con Chile en break-even, el margen "
            "consolidado pasaría de 15,6% a ~31% sin mover nada más.", accent=RED)
    source_note(s, "Fuente 2024: Gastos a Diciembre 2024_v0.xlsx  ·  Fuente 2025: 202601_FLUJO REDCO "
                   "2025 - Cierre diciembre.xlsx (Tabla datos unidades_Dic), vía pilar_a. "
                   "Margen = Ingreso − Costo Op.; Utilidad = Margen ÷ Ingreso. Valores en MUSD.")


def s08_rentabilidad_nivel(prs, d):
    s = blank(prs)
    header(s, "Los estudios largos e intensivos no recuperan su costo",
           "Rentabilidad estimada por nivel de estudio: ingreso cobrado menos costo modelado.")

    niveles = d["niveles"]
    y0 = 1.66
    row_h = 0.475
    max_ing = max(n["ingreso"] for n in niveles)
    for i, n in enumerate(niveles):
        y = y0 + i * row_h
        m = n["margen"]
        col = RED if m < 0.10 else (AMBER if m < 0.25 else GREEN)
        txt(s, MARGIN, y + 0.08, 2.85, 0.26, n["nivel"], size=11, bold=True, color=INK)
        txt(s, MARGIN, y + 0.28, 2.85, 0.18, f"{int(n['n'])} proyectos", size=8.5, color=GREY)
        bar_x, bar_w = 3.50, 5.05
        rect(s, bar_x, y + 0.10, bar_w, 0.26, fill=WHITE, line=LINE)
        w = max(0.04, bar_w * n["ingreso"] / max_ing)
        rect(s, bar_x, y + 0.10, w, 0.26, fill=STEEL)
        txt(s, bar_x + w + 0.08, y + 0.12, 1.2, 0.22,
            f"{fmt_musd(n['ingreso'], 2)} M", size=9, color=GREY)
        rect(s, 9.90, y + 0.06, 1.30, 0.34, fill=BG, line=col)
        txt(s, 9.90, y + 0.11, 1.30, 0.26, pct(m, 1), size=11.5, bold=True,
            color=col, align=PP_ALIGN.CENTER)
        rect(s, 11.36, y + 0.10, 1.45, 0.26, fill=None, line=None)
        txt(s, 11.36, y + 0.12, 1.45, 0.24,
            ("Core" if n["nivel"] not in ("FS e Ing. Basica", "Permisos / MEIA")
             else "Adyacencia"), size=9, color=GREY)

    txt(s, MARGIN, 1.40, 2.85, 0.20, "NIVEL DE ESTUDIO", size=8, bold=True, color=GREY)
    txt(s, 3.50, 1.40, 5.05, 0.20, "INGRESO COBRADO", size=8, bold=True, color=GREY)
    txt(s, 9.90, 1.40, 1.30, 0.20, "MARGEN", size=8, bold=True, color=GREY, align=PP_ALIGN.CENTER)
    txt(s, 11.36, 1.40, 1.45, 0.20, "CLASIFICACIÓN", size=8, bold=True, color=GREY)

    insight(s, MARGIN, 5.60, SLIDE_W - 2 * MARGIN, 1.30, "INSIGHT CLAVE",
            "5,8% vs 67%",
            "Scoping (5,8%) y Confiabilidad de Planes (−8,0%) no recuperan su costo, y entre ambos "
            "explican 1,3 MUSD de ingreso. En el otro extremo, Permisos/MEIA (67%) y Consultoría/Due "
            "Diligence (56%) son los de mayor margen. La cartera rinde 32,9% en conjunto: el problema "
            "no es el precio promedio, es el mix.", accent=AMBER)
    source_note(s, "Fuente: 202601_Proyeccion_Ing+Modulos 1_Edu.xlsx (Gantt + Listas) y "
                   "00_Ciclo Mensual EdP 2023_v0.xlsx (Ingresado), vía pilar_a hoja 06. "
                   "32 proyectos con código mapeado. Ingreso = cobrado real; costo = modelado "
                   "(módulos × costo/módulo × meses), tarifas estándar, no nómina real.")


# ---------------------------------------------------- mercado (ETC / Woodmac)

# Cifras publicadas — ver pie de fuente de cada slide.
ETC_MATERIALES = [
    ("Litio", "Baterías EV y almacenamiento"),
    ("Níquel", "Baterías y acero inoxidable"),
    ("Grafito", "Ánodos de batería"),
    ("Cobalto", "Baterías NMC"),
    ("Neodimio", "Imanes: EV y eólica"),
    ("Cobre", "Redes, EV y electrificación"),
]
CU_ANIOS = list(range(2025, 2036))
CU_DEMANDA_2035, CU_DEMANDA_2025 = 42.7, 42.7 / 1.24
CU_BASE_2035 = CU_DEMANDA_2035 - 8.0 - 3.5   # 31,2 Mt según aritmética Wood Mackenzie
CU_CAPEX_REQ, CU_CAPEX_CHINA = 210.0, 76.0


def _interp(a, b, n):
    return [a + (b - a) * i / (n - 1) for i in range(n)]


def s09_mercado_etc(prs):
    s = blank(prs)
    header(s, "El mercado es promisorio: no falta mineral, falta capacidad de producirlo",
           "Dos lecturas del mismo estudio: cuánto material exige la transición y en qué "
           "materiales el suministro no alcanza a escalar.")

    txt(s, MARGIN, 1.58, 5.9, 0.26, "Necesidad de materiales se incrementa hacia 2030+",
        size=13, bold=True, color=INK)
    line_h(s, MARGIN, 1.90, 5.9, color=LINE)

    rect(s, MARGIN, 2.06, 5.9, 1.36, fill=NAVY)
    txt(s, MARGIN + 0.28, 2.24, 5.35, 0.52, "6.500 millones de toneladas",
        size=27, bold=True, color=WHITE)
    txt(s, MARGIN + 0.28, 2.80, 5.35, 0.50,
        "de materiales de uso final requiere la transición energética entre 2022 y 2050 "
        "para turbinas eólicas, paneles solares y vehículos eléctricos.",
        size=10.5, color=ICE)

    for i, (label, valor, nota) in enumerate([
            ("CONCENTRACIÓN", "95%", "es acero, cobre y aluminio"),
            ("ESCALA COMPARADA", "8.000 Mt", "de carbón se extraen hoy cada año"),
            ("DISPONIBILIDAD", "Sin escasez", "los recursos geológicos superan la demanda 22–50")]):
        y = 3.56 + i * 0.62
        rect(s, MARGIN, y, 5.9, 0.54, fill=BG, line=LINE)
        txt(s, MARGIN + 0.16, y + 0.06, 1.75, 0.20, label, size=8, bold=True, color=GREY)
        txt(s, MARGIN + 0.16, y + 0.24, 1.75, 0.26, valor, size=13, bold=True, color=NAVY)
        txt(s, MARGIN + 2.05, y + 0.15, 3.70, 0.34, nota, size=10, color=INK)

    txt(s, 6.92, 1.58, 5.91, 0.26, "Dónde el suministro no alcanza a escalar",
        size=13, bold=True, color=INK)
    line_h(s, 6.92, 1.90, 5.91, color=LINE)
    txt(s, 6.92, 1.96, 5.91, 0.24,
        "Seis materiales con riesgo de brecha de suministro si no se actúa sobre "
        "eficiencia, reciclaje y oferta minada:", size=10, italic=True, color=GREY)

    for i, (mat, driver) in enumerate(ETC_MATERIALES):
        col, row = i % 2, i // 2
        x = 6.92 + col * 3.02
        y = 2.42 + row * 0.60
        crit = mat in ("Litio", "Cobre")
        rect(s, x, y, 2.86, 0.52, fill=BG, line=RED if crit else LINE)
        oval(s, x + 0.14, y + 0.16, 0.20, fill=RED if crit else STEEL)
        txt(s, x + 0.44, y + 0.06, 2.30, 0.22, mat, size=11.5, bold=True, color=INK)
        txt(s, x + 0.44, y + 0.27, 2.30, 0.20, driver, size=8.5, color=GREY)
    txt(s, 6.92, 4.26, 5.91, 0.22,
        "En rojo: litio y cobre — los más difíciles de escalar en la próxima década.",
        size=9, italic=True, color=RED)

    rect(s, 6.92, 4.60, 5.91, 0.86, fill=BG, line=LINE)
    rect(s, 6.92, 4.60, 0.05, 0.86, fill=STEEL)
    txt(s, 7.12, 4.68, 5.55, 0.20, "Y EUROPA PERDIÓ SU BASE PRODUCTIVA", size=8, bold=True, color=GREY)
    txt(s, 7.12, 4.88, 5.55, 0.50,
        "La participación europea en la producción mundial de minerales cayó de 25% a menos "
        "de 7% en 40 años — el mismo vacío que empuja el friend-shoring hacia las Américas.",
        size=10, color=INK)

    rect(s, MARGIN, 5.60, SLIDE_W - 2 * MARGIN, 1.06, fill=NAVY)
    txt(s, MARGIN + 0.28, 5.72, 11.8, 0.34,
        "REDCO tiene y tendrá la oportunidad de capturar la necesidad de pasar de Recursos a "
        "Reservas", size=15, bold=True, color=WHITE)
    txt(s, MARGIN + 0.28, 6.12, 11.8, 0.42,
        "El cuello de botella no es geológico sino de ejecución: estudios, ingeniería y "
        "permisos. Ese es exactamente el CORE de REDCO.", size=11, color=ICE)
    source_note(s, "Fuente: Energy Transitions Commission, «Material and Resource Requirements for "
                   "the Energy Transition» (julio 2023) y «A Critical Raw Material Supply-Side "
                   "Innovation Roadmap for the EU Energy Transition» (diciembre 2024) — "
                   "energy-transitions.org. Misma fuente citada en el deck 2024, en su versión vigente.")


def _slide_cobre(prs, variante):
    s = blank(prs)
    header(s, f"En cobre faltan más de {CU_CAPEX_REQ:.0f} bUS$ para cubrir la demanda al 2035",
           "La capacidad instalada y planificada no alcanza: la brecha se cierra con nueva "
           "mina y con chatarra.")

    txt(s, MARGIN, 1.56, 7.15, 0.24, "Oferta vs. demanda de cobre  ·  Mt", size=12,
        bold=True, color=INK)
    n = len(CU_ANIOS)
    base = _interp(CU_DEMANDA_2025, CU_BASE_2035, n)
    nueva = _interp(0.0, 8.0, n)
    chatarra = _interp(0.0, 3.5, n)
    demanda = _interp(CU_DEMANDA_2025, CU_DEMANDA_2035, n)
    ch = add_chart(s, XL_CHART_TYPE.AREA_STACKED, MARGIN - 0.05, 1.84, 7.30, 3.12,
                   [str(a) for a in CU_ANIOS],
                   [("Capacidad existente y planificada", base),
                    ("Nueva capacidad requerida", nueva),
                    ("Chatarra directa", chatarra)],
                   colors=[NAVY, MID, SKY], legend=True, font_size=8.5, y_min=25)
    ch.plots[0].vary_by_categories = False

    rect(s, 1.05, 1.92, 3.05, 0.66, fill=WHITE, line=STEEL)
    txt(s, 1.17, 1.98, 2.85, 0.20, "BRECHA A 2035", size=8, bold=True, color=GREY)
    txt(s, 1.17, 2.16, 2.85, 0.34, "8,0 Mt nueva mina + 3,5 Mt chatarra",
        size=10.5, bold=True, color=NAVY)
    txt(s, MARGIN, 5.04, 7.30, 0.22,
        "El tope de la pila es la demanda total: 34,4 Mt (2025) → 42,7 Mt (2035).",
        size=8.5, italic=True, color=GREY)

    txt(s, 8.10, 1.56, 4.73, 0.24, "Inversión de capital  ·  bUS$", size=12, bold=True, color=INK)
    add_chart(s, XL_CHART_TYPE.COLUMN_CLUSTERED, 8.05, 1.84, 4.78, 3.12,
              ["Requerido\na 2035", "Referencia:\nfinanciamiento\nchino 2019–25"],
              [("bUS$", [CU_CAPEX_REQ, CU_CAPEX_CHINA])],
              colors=[NAVY], legend=False, gap=110, font_size=8.5)
    ch2 = s.shapes[-1].chart
    ch2.plots[0].has_data_labels = True
    dl2 = ch2.plots[0].data_labels
    dl2.font.size = Pt(15); dl2.font.bold = True; dl2.font.color.rgb = NAVY
    dl2.position = XL_LABEL_POSITION.OUTSIDE_END

    for i, (label, valor) in enumerate([
            ("DEMANDA 2035", "42,7 Mt  ·  +24%"),
            ("COBERTURA DE MINAS ACTUALES", "~70% de la demanda 2035"),
            ("DÉFICIT REFINADO 2025", "304 kt, mayor en 2026")]):
        x = MARGIN + i * 4.28
        rect(s, x, 5.36, 4.06, 0.52, fill=BG, line=LINE)
        txt(s, x + 0.14, 5.42, 3.80, 0.18, label, size=8, bold=True, color=GREY)
        txt(s, x + 0.14, 5.60, 3.80, 0.24, valor, size=11.5, bold=True, color=INK)

    if variante == "A":
        titulo = "Si consideramos el 10% en gastos de ingeniería ≈ 21 bUS$ de mercado potencial"
        detalle = ("Mismo criterio que el deck 2024 (10% del CAPEX). Simple y comparable año "
                   "contra año, pero es un supuesto propio: no proviene de Wood Mackenzie ni de IEA.")
        chip = "VERSIÓN A · SUPUESTO 10%"
    else:
        titulo = "Con el 4–7% de ingeniería sobre CAPEX ≈ 8,4 a 14,7 bUS$ de mercado"
        detalle = ("Se aplica el rango θ de ingeniería+consultoría sobre CAPEX minero usado en el "
                   "informe de demanda REDCO. Se presenta como rango, no como punto: es más "
                   "defendible ante el cliente que un porcentaje único.")
        chip = "VERSIÓN B · FRACCIÓN DE INGENIERÍA 4–7%"

    rect(s, MARGIN, 5.96, SLIDE_W - 2 * MARGIN, 1.00, fill=NAVY)
    txt(s, MARGIN + 0.26, 6.04, 8.55, 0.34, titulo, size=13.5, bold=True, color=WHITE)
    txt(s, MARGIN + 0.26, 6.42, 8.55, 0.46, detalle, size=9.5, color=ICE)
    rect(s, 9.55, 6.28, 3.22, 0.32, fill=STEEL)
    txt(s, 9.55, 6.34, 3.22, 0.22, chip, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    source_note(s, "Fuente: Wood Mackenzie, «High-wire act: is soaring copper demand an obstacle to "
                   "future growth?» (octubre 2025) — demanda 42,7 Mtpa a 2035 (+24%), 8 Mtpa de nueva "
                   "capacidad, 3,5 Mtpa de chatarra, >210 bUS$ de inversión, 76 bUS$ de financiamiento "
                   "chino 2019-25. IEA, «Global Critical Minerals Outlook» — minas existentes y "
                   "planificadas cubren ~70% de la demanda 2035. Los puntos 2025 y 2035 son de fuente; "
                   "la trayectoria intermedia es interpolación lineal.")
    return s


def s10a_cobre(prs):
    _slide_cobre(prs, "A")


def s10b_cobre(prs):
    _slide_cobre(prs, "B")


def s11_fundamentos(prs):
    s = blank(prs)
    header(s, "Los principales fundamentos de la estrategia",
           "Cinco fundamentos definidos en 2024; la data 2025–2026 confirma cuatro y obliga a "
           "revisar uno.")
    fundamentos = [
        ("Nuevas tendencias de mercado garantizan una creciente demanda de minerales de "
         "nuestros clientes", "Confirmado", GREEN,
         "ETC: 6.500 Mt de materiales a 2050; brecha en litio y cobre."),
        ("La Transición Energética requerirá a nuestros clientes nuevos estándares "
         "ambientales y sociales", "Confirmado", GREEN,
         "Permisos/MEIA es el nivel de mayor margen de la cartera (67%)."),
        ("El contexto está demandando a nuestros clientes y generará un ciclo de escasez "
         "de profesionales", "Confirmado", GREEN,
         "La dotación creció 40% en un año (81 → 113 personas)."),
        ("Ejecución de calidad de proyectos es esencial para sólidas relaciones con "
         "clientes clave", "Confirmado", GREEN,
         "Los clientes efectivos pasaron de 6,6 a 18,5 sin perder los históricos."),
        ("Financieramente es una compañía sana y con capacidad de autosoportar su "
         "crecimiento", "Requiere revisión", RED,
         "El margen cayó de 26,5% a 13,5% y Chile opera a −39%."),
    ]
    y = 1.62
    for i, (texto, estado, color, evidencia) in enumerate(fundamentos, 1):
        h = 0.94
        rect(s, MARGIN, y, SLIDE_W - 2 * MARGIN, h, fill=BG, line=LINE)
        o = oval(s, MARGIN + 0.20, y + 0.31, 0.32, fill=STEEL)
        shape_text(o, str(i), size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER, pad=0)
        txt(s, MARGIN + 0.66, y + 0.14, 7.15, 0.44, texto, size=12.5, bold=True, color=INK)
        txt(s, MARGIN + 0.66, y + 0.60, 7.15, 0.26, evidencia, size=9.5, italic=True, color=GREY)
        rect(s, 9.40, y + 0.28, 1.55, 0.36, fill=WHITE, line=color)
        txt(s, 9.40, y + 0.34, 1.55, 0.24, estado, size=9.5, bold=True, color=color,
            align=PP_ALIGN.CENTER)
        txt(s, 11.15, y + 0.34, 1.68, 0.24,
            "Evidencia pilar_a" if i > 2 else "Evidencia mercado", size=9, color=GREY)
        y += h + 0.12
    source_note(s, "Texto de los cinco fundamentos: deck de Gestión Estratégica REDCO, junio 2024 "
                   "(diapositiva 21), conservado literal. La columna de estado es contraste contra "
                   "la data 2025–2026 de pilar_a y las fuentes de mercado de este capítulo.")


def s12_core(prs):
    s = blank(prs)
    header(s, "Para crecer con rentabilidad sostenible, lo primordial es cuidar el CORE",
           "El CORE es el activo y las capacidades más distintivas y potentes — y desde ahí se "
           "desarrollan las adyacencias.")
    rect(s, MARGIN, 1.62, SLIDE_W - 2 * MARGIN, 0.74, fill=NAVY)
    txt(s, MARGIN + 0.30, 1.78, 12.0, 0.42,
        "El CORE es el activo y las capacidades más distintivas y potentes",
        size=17, bold=True, color=WHITE)
    cols = [
        ("Por qué el CORE define dónde ganas", [
            "Define dónde y por qué puedes ganar y generar utilidades.",
            "Es la razón por la que tus clientes más leales y rentables te eligen.",
            "Te da una mirada correcta de quiénes son tus verdaderos competidores."]),
        ("Por qué el CORE te protege y te habilita", [
            "Es tu refugio en condiciones adversas de mercado.",
            "Habilita tus mejores oportunidades de crecimiento.",
            "Da claridad para asignar recursos y transparencia en las decisiones."]),
    ]
    for i, (titulo, items) in enumerate(cols):
        x = MARGIN + i * 6.32
        card(s, x, 2.66, 6.02, 2.42, i + 1, titulo, items)
    insight(s, MARGIN, 5.28, SLIDE_W - 2 * MARGIN, 1.52, "CÓMO SE APLICA ESTO EN REDCO", None,
            "La cartera 2024–25 muestra el CORE (FEL 0-1-2A-2B) en 9,23 MUSD emitidos, un 52% del "
            "total, y la adyacencia (FEL 3-4) en 2,08 MUSD, un 12%. El 33% restante entra sin fase "
            "asignada. Dentro del CORE conviven los niveles de mayor margen (Permisos 67%, "
            "Consultoría 56%, Conceptual/LOM 41%) y los de menor (Scoping 5,8%): cuidar el CORE no "
            "es solo defender la fase, es corregir el mix dentro de ella.", accent=NAVY)
    source_note(s, "Texto conservado del deck de Gestión Estratégica REDCO, junio 2024 "
                   "(diapositiva 22). Cifras de cartera: " + FUENTE_PILAR_A)


# ================================================================== CAPÍTULO 3

EJES_CORE = [
    ("Países", 0), ("Alianzas", 55), ("Tipo de cliente", 125),
    ("Ámbitos de acción", 180), ("Especialidades\nTécnicas", 235), ("Tipo de producto", 305),
]

CORE_TEXTO = ("Consultora boutique de confianza, reconocida por su eficiencia y flexibilidad. "
              "Destacados en planificación y desarrollo minero en Chile, Perú y Brasil.")


def radial_core(s, cx, cy, r, movimientos=(), destacar=()):
    """Diagrama radial CORE / adyacencias: anillos concéntricos + 6 ejes."""
    for k in (1.0, 0.78, 0.56):
        rr = r * k
        rect(s, cx - rr, cy - rr, rr * 2, rr * 2, fill=None, line=LINE,
             shape=MSO_SHAPE.OVAL)
    rect(s, cx - r * 0.40, cy - r * 0.40, r * 0.80, r * 0.80, fill=WHITE, line=NAVY,
         line_w=1.5, shape=MSO_SHAPE.OVAL)
    if r >= 1.75:
        txt(s, cx - r * 0.34, cy - r * 0.26, r * 0.68, r * 0.55, CORE_TEXTO,
            size=6.5, color=NAVY, align=PP_ALIGN.CENTER)
    else:
        # con radio chico el párrafo se sale del círculo: basta el rótulo
        txt(s, cx - r * 0.34, cy - 0.12, r * 0.68, 0.26, "CORE",
            size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    for label, ang in EJES_CORE:
        rad = math.radians(ang)
        ex, ey = cx + math.cos(rad) * r * 1.13, cy - math.sin(rad) * r * 1.13
        conn = s.shapes.add_connector(1, Inches(cx), Inches(cy), Inches(ex), Inches(ey))
        conn.line.color.rgb = STEEL
        conn.line.width = Pt(0.75)
        lx, ly = cx + math.cos(rad) * r * 1.24, cy - math.sin(rad) * r * 1.24
        activo = label.replace("\n", " ") in destacar
        align = PP_ALIGN.CENTER
        txt(s, lx - 0.72, ly - 0.14, 1.44, 0.34, label.replace("\n", " "), size=8.5,
            bold=activo, color=(INK if activo else GREY), align=align)

    for num, ang, frac in movimientos:
        rad = math.radians(ang)
        mx, my = cx + math.cos(rad) * r * frac, cy - math.sin(rad) * r * frac
        o = oval(s, mx - 0.125, my - 0.125, 0.25, fill=WHITE, line=RED)
        shape_text(o, num, size=8, bold=True, color=RED, align=PP_ALIGN.CENTER, pad=0)


def harvey(s, x, y, d, fraccion, color=STEEL):
    """Harvey ball: círculo con relleno parcial que indica cercanía al CORE.

    El sector se dibuja como freeform y no con la geometría `pie`: los presets no
    respetan ángulos negativos de forma fiable entre PowerPoint y LibreOffice, y
    un 2% terminaba renderizándose como círculo completo.
    """
    r = d / 2.0
    cx, cy = x + r, y + r
    if fraccion >= 0.995:
        rect(s, x, y, d, d, fill=color, line=GREY, line_w=0.75, shape=MSO_SHAPE.OVAL)
        return
    rect(s, x, y, d, d, fill=WHITE, line=GREY, line_w=0.75, shape=MSO_SHAPE.OVAL)
    if fraccion > 0.005:
        pasos = max(3, int(round(48 * fraccion)))
        puntos = []
        for i in range(pasos + 1):
            ang = math.radians(-90 + 360.0 * fraccion * i / pasos)
            puntos.append((Inches(cx + r * math.cos(ang)), Inches(cy + r * math.sin(ang))))
        ff = s.shapes.build_freeform(Inches(cx), Inches(cy))
        ff.add_line_segments(puntos, close=True)
        wedge = ff.convert_to_shape()
        wedge.fill.solid()
        wedge.fill.fore_color.rgb = color
        wedge.line.fill.background()
        wedge.shadow.inherit = False
    # contorno limpio por encima del sector
    rect(s, x, y, d, d, fill=None, line=GREY, line_w=0.75, shape=MSO_SHAPE.OVAL)


def s13_divisor_cap3(prs):
    divider(blank(prs), "03", "Metodología para la toma de decisiones",
            "Cómo se elige dónde crecer: atractivo del mercado, habilidad de ganar y disciplina "
            "de distancia al CORE.")


def s14_elementos(prs):
    s = blank(prs)
    header(s, "Las mejores estrategias aplican una estricta revisión de estos elementos",
           "Dos dimensiones de revisión: cuán atractivo es el mercado y cuán capaces somos "
           "de ganar en él.")
    columnas = [
        ("Atractivo del mercado", MARGIN, [
            ("Apuntar a posibilidades futuras grandes y con expectativa de crecimiento rentable",
             ""),
            ("Moverse por las corrientes digitales y de sostenibilidad",
             "Cambian los límites comerciales y crean nuevas palancas de ventaja competitiva.")]),
        ("Habilidad de ganar", 6.92, [
            ("Construir sobre las fortalezas principales",
             "Aprovechar los activos, las capacidades diferenciadoras y las relaciones "
             "íntimas con los clientes."),
            ("Crear un camino para ser líderes de manera rentable",
             "Muchos seguidores en mercados atractivos se suman a una buena estrategia."),
            ("Construir con una fórmula repetible",
             "Crear una plataforma de crecimiento, invirtiendo con múltiples movimientos y "
             "construyendo un negocio a escala.")]),
    ]
    TOP, BOT, GAPY = 2.04, 5.18, 0.14
    for titulo, x, items in columnas:
        w = 5.91
        txt(s, x, 1.58, w, 0.28, titulo, size=14, bold=True, color=INK, align=PP_ALIGN.CENTER)
        line_h(s, x, 1.92, w, color=LINE)
        n = len(items)
        h = (BOT - TOP - GAPY * (n - 1)) / n
        y = TOP
        for i, (cab, det) in enumerate(items, 1):
            rect(s, x, y, w, h, fill=BG, line=LINE)
            o = oval(s, x + 0.18, y + 0.16, 0.28, fill=STEEL)
            shape_text(o, str(i), size=10.5, bold=True, color=WHITE,
                       align=PP_ALIGN.CENTER, pad=0)
            txt(s, x + 0.56, y + 0.14, w - 0.76, 0.44, cab, size=12, bold=True, color=INK)
            if det:
                txt(s, x + 0.56, y + h - 0.56, w - 0.76, 0.48, det, size=10, color=GREY)
            y += h + GAPY

    rect(s, MARGIN, 5.34, SLIDE_W - 2 * MARGIN, 1.50, fill=NAVY)
    txt(s, MARGIN + 0.28, 5.46, 5.55, 0.80,
        "El atractivo de REDCO está en los clientes y países donde su CORE tiene la mayor "
        "capacidad de generar ventaja competitiva en ciclos sostenibles.",
        size=11.5, bold=True, color=WHITE)
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.28), Inches(5.86),
                           Inches(0.52), Inches(0.42))
    a.fill.solid(); a.fill.fore_color.rgb = ICE; a.line.fill.background()
    a.shadow.inherit = False
    txt(s, 7.02, 5.46, 5.55, 1.24,
        [("•  REDCO apalanca desde las especialidades diferenciadoras en planificación y "
          "desarrollo minero.", {"size": 10, "spacing": 4}),
         ("•  Genera relaciones desde la confianza, la eficiencia y la flexibilidad, con "
          "estructura de costos rentable.", {"size": 10, "spacing": 4}),
         ("•  En la reiteración es capaz de generar una plataforma sistemática de "
          "crecimiento rentable.", {"size": 10})], color=ICE)
    source_note(s, "Contenido conservado del deck de Gestión Estratégica REDCO, junio 2024 "
                   "(diapositiva 31). Rediseñado bajo el Estándar de Presentaciones REDCO v1.0.")


def s15_flujo_decisiones(prs):
    s = blank(prs)
    header(s, "El flujo de decisiones para crecer y cuidar el CORE",
           "Cuatro movimientos ordenados: primero país, después producto, después "
           "reconocimiento, y por último tipo de cliente.")
    rect(s, MARGIN, 1.56, 7.55, 5.30, fill=BG, line=LINE)
    radial_core(s, MARGIN + 3.78, 4.20, 2.05,
                movimientos=[("1", 8, 0.92), ("2", 318, 0.80),
                             ("3", 232, 0.86), ("4", 128, 0.96),
                             ("1.1\n4.1", 58, 0.72)],
                destacar=("Países", "Tipo de cliente", "Tipo de producto"))
    pasos = [
        ("1", "Ingresar a un país nuevo con un cliente minero",
         "1.1  Evaluar las alianzas necesarias"),
        ("2", "Crecer en ese país con el set de productos del CORE de REDCO", ""),
        ("3", "Fortalecer el set de productos en ese país hasta ser boutique reconocida",
         "Fortalecer la experiencia del equipo para ese país"),
        ("4", "Crecer en cliente más allá de los mineros",
         "4.1  Evaluar las alianzas necesarias"),
    ]
    txt(s, 8.42, 1.58, 4.41, 0.52,
        "Hipótesis de la dinámica estratégica entre CORE y adyacencia",
        size=13, bold=True, color=INK)
    line_h(s, 8.42, 2.16, 4.41, color=LINE)
    y = 2.30
    for num, cab, sub in pasos:
        h = 1.10 if sub else 0.80
        rect(s, 8.42, y, 4.41, h, fill=BG, line=LINE)
        o = oval(s, 8.60, y + 0.18, 0.30, fill=STEEL)
        shape_text(o, num, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER, pad=0)
        txt(s, 9.00, y + 0.14, 3.68, 0.56, cab, size=11, bold=True, color=INK)
        if sub:
            txt(s, 9.00, y + 0.72, 3.68, 0.30, sub, size=9.5, italic=True, color=GREY)
        y += h + 0.11
    source_note(s, "Contenido conservado del deck de Gestión Estratégica REDCO, junio 2024 "
                   "(diapositiva 32). El diagrama se reconstruyó como formas nativas.")


CRITERIOS_CORE = [
    "Tipo de clientes conocidos",
    "Estructura de costos conocidos",
    "Ámbitos de acción conocidos",
    "Especialidades técnicas conocidas",
    "Tipo de competidores conocidos",
]
PASOS_ADY = [
    ("CORE", "—", [1.0, 1.0, 1.0, 1.0, 1.0]),
    ("1 paso de la adyacencia", "35%", [0.95, 0.60, 0.55, 0.62, 0.42]),
    ("2 pasos de la adyacencia", "15%", [0.50, 0.45, 0.42, 0.72, 0.45]),
    ("3 pasos de la adyacencia", "8%", [0.05, 0.45, 0.48, 0.45, 0.05]),
    ("Diversificación", "-5%", [0.02, 0.02, 0.02, 0.02, 0.02]),
]


def _matriz_core(s, x0, y0, pasos, col_w=1.85, row_h=0.62, label_w=3.05):
    txt(s, x0, y0 - 0.52, label_w, 0.24, "RELACIÓN CON EL CORE", size=8.5, bold=True, color=GREY)
    txt(s, x0, y0 - 0.28, label_w, 0.24, "Probabilidades de éxito", size=10.5,
        italic=True, color=GREY)
    for j, (titulo, prob, _) in enumerate(pasos):
        cx = x0 + label_w + j * col_w
        if titulo == "CORE":
            rect(s, cx + 0.16, y0 - 0.62, col_w - 0.34, 0.42, fill=STEEL)
            txt(s, cx + 0.16, y0 - 0.55, col_w - 0.34, 0.26, "CORE", size=10.5, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)
        else:
            txt(s, cx, y0 - 0.62, col_w - 0.10, 0.40, titulo, size=9.5, color=INK,
                align=PP_ALIGN.CENTER)
        if prob != "—":
            col = GREEN if prob == "35%" else (AMBER if prob == "15%" else RED)
            txt(s, cx, y0 - 0.24, col_w - 0.10, 0.24, prob, size=12, bold=True, color=col,
                align=PP_ALIGN.CENTER)
    for i, crit in enumerate(CRITERIOS_CORE):
        y = y0 + 0.10 + i * row_h
        o = oval(s, x0, y + 0.02, 0.22, fill=WHITE, line=RED)
        shape_text(o, str(i + 1), size=8, bold=True, color=RED, align=PP_ALIGN.CENTER, pad=0)
        txt(s, x0 + 0.34, y, label_w - 0.44, 0.40, crit, size=10, color=INK)
        for j, (_, _, fracs) in enumerate(pasos):
            cx = x0 + label_w + j * col_w + (col_w - 0.10) / 2 - 0.15
            harvey(s, cx, y - 0.02, 0.30, fracs[i],
                   color=STEEL if j == 0 else RGB_DARK)


RGB_DARK = GREY


def s16_distancia_core(prs):
    s = blank(prs)
    header(s, "Ser disciplinados y mantenerse conectados con el CORE es crucial para el éxito",
           "A mayor distancia económica desde el CORE, menos atributos conocidos y menor "
           "probabilidad de éxito.")
    _matriz_core(s, MARGIN, 2.30, PASOS_ADY)
    rect(s, MARGIN, 5.72, SLIDE_W - 2 * MARGIN, 1.02, fill=NAVY)
    txt(s, MARGIN + 0.28, 5.86, 4.0, 0.30, "Distancia económica desde el CORE",
        size=13, bold=True, color=WHITE)
    txt(s, 5.40, 5.82, 7.4, 0.84,
        [("•  Menor ventaja competitiva", {"size": 10, "spacing": 2}),
         ("•  Menor potencial de generar valor", {"size": 10, "spacing": 2}),
         ("•  Menor participación y lealtad con clientes", {"size": 10})], color=ICE)
    source_note(s, "Contenido conservado del deck de Gestión Estratégica REDCO, junio 2024 "
                   "(diapositiva 33). Probabilidades de éxito por grado de adyacencia según el "
                   "marco de Bain & Company sobre crecimiento adyacente al núcleo.")


PASOS_USA = [("CORE", "—", [1.0, 1.0, 1.0, 1.0, 1.0]),
             ("USA · 1 paso de la adyacencia", "35%", [0.95, 0.60, 0.55, 0.62, 0.42])]

USA_LECTURA = [
    "Entramos con clientes mineros, el tipo de cliente que ya conocemos.",
    "Mantenemos la estructura de costos conocida de Chile.",
    "Entramos con proyectos de planificación minera, nuestro ámbito de acción.",
    "Entramos con proyectos en minería, nuestra especialidad técnica.",
    "Percibimos espacio para ganar frente a los competidores presentes en licitaciones.",
]


def s17_usa(prs, d):
    s = blank(prs)
    header(s, "Aplicado a USA: la entrada se hizo a un solo paso de adyacencia",
           "Los cinco criterios del CORE se mantuvieron conocidos, y el resultado lo confirma.")
    _matriz_core(s, MARGIN, 2.30, PASOS_USA, col_w=2.10, label_w=2.85)
    txt(s, 7.60, 1.68, 5.23, 0.26, "Por qué USA calificó como un paso", size=13,
        bold=True, color=INK)
    line_h(s, 7.60, 2.00, 5.23, color=LINE)
    for i, linea in enumerate(USA_LECTURA):
        y = 2.14 + i * 0.62
        rect(s, 7.60, y, 5.23, 0.54, fill=BG, line=LINE)
        o = oval(s, 7.76, y + 0.16, 0.22, fill=STEEL)
        shape_text(o, str(i + 1), size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER, pad=0)
        txt(s, 8.10, y + 0.08, 4.60, 0.40, linea, size=10, color=INK)
    v = d["venta_pais"]["USA"]
    insight(s, MARGIN, 5.44, SLIDE_W - 2 * MARGIN, 1.36, "LA HIPÓTESIS SE VALIDÓ",
            "0,10 → 1,65 MUSD",
            "USA pasó de 0,10 MUSD en 2023 a 1,65 MUSD en 2025 y opera con 63% de utilidad sobre "
            "ingreso, la segunda más alta del grupo. Con 3 personas de dotación en 2026, es la "
            "operación más eficiente por persona. La disciplina de entrar a un solo paso del CORE "
            "es exactamente lo que hizo replicable el movimiento.", accent=GREEN)
    source_note(s, "Marco conservado del deck 2024 (diapositiva 34). Cifras de validación: "
                   + FUENTE_PILAR_A)


def s18_probabilidades(prs):
    s = blank(prs)
    header(s, "Las probabilidades de éxito son altas creciendo en países y tipo de clientes",
           "Los dos movimientos elegidos se mantienen a un paso del CORE; producto, "
           "especialidad y ámbito no se expanden.")
    _matriz_core(s, MARGIN, 2.26, PASOS_USA, col_w=1.95, label_w=2.70)
    rect(s, 7.30, 1.62, 5.53, 4.16, fill=BG, line=LINE)
    radial_core(s, 10.06, 3.60, 1.34,
                movimientos=[("1", 8, 0.92), ("2", 318, 0.80),
                             ("3", 232, 0.86), ("4", 128, 0.96)],
                destacar=("Países", "Tipo de cliente"))
    txt(s, 7.46, 5.32, 5.20, 0.34,
        "Se crece en países y en tipo de cliente; las alianzas son el habilitador.",
        size=10, italic=True, color=GREY, align=PP_ALIGN.CENTER)
    decisiones = [
        ("Se crece", "Países", NAVY), ("Se crece", "Tipo de cliente", NAVY),
        ("Habilitador", "Alianzas", STEEL), ("No se expande", "Tipo de producto", SLATE),
        ("No se expande", "Especialidades técnicas", SLATE),
        ("No se expande", "Ámbitos de acción", SLATE),
    ]
    for i, (estado, eje, col) in enumerate(decisiones):
        x = MARGIN + (i % 3) * 4.28
        y = 5.94 + (i // 3) * 0.50
        rect(s, x, y, 4.06, 0.42, fill=WHITE, line=LINE)
        rect(s, x, y, 0.05, 0.42, fill=col)
        txt(s, x + 0.16, y + 0.05, 1.30, 0.20, estado.upper(), size=7.5, bold=True, color=col)
        txt(s, x + 0.16, y + 0.21, 3.70, 0.20, eje, size=10, bold=True, color=INK)
    source_note(s, "Contenido conservado del deck de Gestión Estratégica REDCO, junio 2024 "
                   "(diapositiva 35). Rediseñado como formas nativas.")


# ================================================================== CAPÍTULO 4

def _flecha(s, x1, y1, x2, y2, color=RED, width=2.25):
    """Conector recto con punta de flecha (python-pptx no expone la punta)."""
    from pptx.oxml.ns import qn
    c = s.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    ln = c.line._get_or_add_ln()
    head = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
    ln.append(head)
    return c


def _motor(s, x, y, w, h, eje_x, eje_y, color):
    """Motivo de crecimiento: ejes + paralelogramo + flecha diagonal."""
    ax, ay = x + 0.34, y + h - 0.42
    c = s.shapes.add_connector(1, Inches(ax), Inches(ay), Inches(x + w - 0.10), Inches(ay))
    c.line.color.rgb = INK; c.line.width = Pt(1)
    c = s.shapes.add_connector(1, Inches(ax), Inches(ay), Inches(ax), Inches(y + 0.20))
    c.line.color.rgb = INK; c.line.width = Pt(1)
    p = s.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, Inches(ax + 0.16), Inches(y + 0.34),
                           Inches(w - 0.72), Inches(h - 0.94))
    p.fill.background()
    p.line.color.rgb = GREY
    p.line.width = Pt(1.25)
    p.shadow.inherit = False
    _flecha(s, ax + 0.30, ay - 0.14, ax + w - 0.72, y + 0.44, color=color)
    txt(s, x + 0.06, y + 0.16, 1.5, 0.20, eje_y, size=8.5, bold=True, color=INK)
    txt(s, x + w - 1.72, y + h - 0.34, 1.66, 0.22, eje_x, size=8.5, bold=True, color=INK,
        align=PP_ALIGN.RIGHT)


def s19_divisor_cap4(prs):
    divider(blank(prs), "04", "Propósito, misión y habilitadores",
            "El propósito y la misión vigentes, los motores de crecimiento y las ocho metas "
            "al horizonte 2028.")


def s20_motores(prs):
    s = blank(prs)
    header(s, "Aplicando estos conceptos al crecimiento de REDCO con mayor probabilidad de ganar",
           "El equilibrio entre ampliar horizontes y consolidar el CORE se descompone en "
           "cuatro palancas de crecimiento.")
    rect(s, MARGIN, 1.62, 5.30, 4.44, fill=BG, line=LINE)
    _motor(s, MARGIN + 0.30, 1.86, 4.70, 3.66,
           "Velocidad de consolidación del CORE", "Velocidad de crecimiento", NAVY)
    txt(s, MARGIN + 0.30, 5.66, 4.70, 0.30, "Ampliar horizontes  ·  consolidar el CORE",
        size=10, italic=True, color=GREY, align=PP_ALIGN.CENTER)

    motores = [
        ("Tipo cliente", "Ventas", STEEL, 6.28, 1.62),
        ("N. Países", "Ventas", STEEL, 9.62, 1.62),
        ("Alianzas", "$ Ventas", STEEL, 6.28, 3.86),
        ("Profesionales senior", "$ Ventas", STEEL, 9.62, 3.86),
    ]
    for eje_x, eje_y, color, x, y in motores:
        rect(s, x, y, 3.22, 2.10, fill=BG, line=LINE)
        _motor(s, x + 0.06, y + 0.06, 3.10, 1.98, eje_x, eje_y, color)

    rect(s, 6.28, 6.14, 6.55, 0.72, fill=NAVY)
    txt(s, 6.46, 6.24, 6.20, 0.52,
        "Cada palanca sube la venta sin alejarse del CORE: por eso la probabilidad de ganar "
        "se mantiene alta.", size=10.5, color=WHITE)
    source_note(s, "Contenido conservado del deck de Gestión Estratégica REDCO, junio 2024 "
                   "(diapositiva 37). Rediseñado como formas nativas.")


METAS_CRECIMIENTO = [
    ("Venta", "9,5 → 14 MUSD al 2028", "CAGR ~14%. Base: 7,1 (’23) · 8,4 (’24) · 9,5 MUSD (’25)."),
    ("Diversificación", "Países efectivos 3,9 → ≥4,5", "Ningún país sobre 40% del emitido "
     "(Chile ya bajó de 61% a 34%)."),
    ("Nuevos mercados", "USA 1,7 → ≥3,0 MUSD", "Y al menos 2 clientes no mineros recurrentes "
     "(precedente: BVI 0,27 MUSD)."),
    ("Autonomía de liderazgo", "1 líder con P&L por país al 2027", "En Perú, Brasil y USA — "
     "reduce la dependencia del dueño-fundador."),
]
METAS_CORE = [
    ("Rentabilidad", "Margen 15,6% → ≥25% al 2028", "Con Chile en break-even operacional al "
     "2027 (hoy −1,41 MUSD)."),
    ("Conversión comercial", "Por valor 23,5% → ≥35%", "Y ticket adjudicado de $132k a ≥$200k "
     "(por número ya está en ~48%)."),
    ("Ciclo de cobro", "60 → ≤45 días", "Rotación de 6,1× a 8×/año; ningún país sobre 60 d "
     "(hoy Brasil 76, Rusia 80)."),
    ("Mix hacia el CORE", "Core FEL 0-2B 52% → ≥60%", "Y los niveles con margen bajo 10% "
     "(Scoping, Confiabilidad) ≤15% del ingreso."),
]

HABILITADORES = [
    ("Fortalecer talento", "Equipos (módulos) autogestionados, de calidad, diversidad y "
     "pertenencia como base al trabajo técnico y administrativo."),
    ("Estandarizar productos", "Estandarizar lo mínimo que garantice calidad y disponibilidad "
     "de experiencias."),
    ("Potenciar la experiencia con el cliente", "Fortalecer de inicio a fin el acompañamiento "
     "para lograr el reconocimiento diferenciador."),
    ("Evolucionar talento, producto y experiencia", "Gestionar los aprendizajes para lograr la "
     "adaptación constante a nuestro entorno."),
]

CADENA_VALOR = [
    ("Comercial", "Lograr calidad/precio mejor que la competencia"),
    ("Ejecución", "En calidad, tiempo y costo, con feedback continuo"),
    ("Gestión de alcance, entregables y cobranza", "Asegurar coherencia entre lo ejecutado y lo comprometido"),
    ("Post evaluación", "Cliente y equipo reconocen una experiencia distintiva"),
]


def _banda(s, y, h, etiqueta):
    txt(s, MARGIN, y + 0.04, 1.06, 0.50, etiqueta, size=8.5, bold=True, color=GREY)
    rect(s, 1.62, y, 0.045, h, fill=STEEL)
    return 1.76, 11.07


def s21_proposito_metas(prs):
    s = blank(prs)
    header(s, "Propósito, misión, metas y habilitadores",
           "Seis niveles encadenados: del propósito a los habilitadores, con ocho metas al 2028 "
           "que hacen medible la aspiración.")

    # 1 · propósito
    y = 1.50
    x, w = _banda(s, y, 0.40, "PROPÓSITO\nEMERGENTE")
    rect(s, x, y, w, 0.40, fill=NAVY)
    txt(s, x + 0.20, y + 0.07, w - 0.40, 0.28,
        "Creemos en un camino de excelencia y aprendizaje creando soluciones para un mundo mejor.",
        size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # 2 · misión
    y = 1.96
    x, w = _banda(s, y, 0.52, "MISIÓN")
    rect(s, x, y, w, 0.52, fill=BG, line=LINE)
    txt(s, x + 0.20, y + 0.06, w - 0.40, 0.42,
        "Somos una consultora boutique que proporciona soluciones a medida e integradas, "
        "logrando una nueva manera de hacer minería, con un gran equipo de profesionales que "
        "maximiza el valor de nuestros clientes.", size=10.5, color=INK, align=PP_ALIGN.CENTER)

    # 3 · motores de crecimiento
    y = 2.54
    x, w = _banda(s, y, 0.84, "MOTORES DE\nCRECIMIENTO")
    prior = [
        ("Primera prioridad · Core", "Equipos ágiles y autónomos, mayor calidad/precio y valor "
         "agregado a la experiencia del cliente.", NAVY),
        ("Segunda prioridad · Adyacencias", "Consolidar y hacer crecer la boutique en Chile, "
         "Brasil y Perú.", STEEL),
        ("Tercera prioridad · Nuevas adyacencias", "Crecer y convertirse en boutique en USA "
         "y Rusia.", GREY),
    ]
    cw = (w - 0.24) / 3
    for i, (titulo, det, col) in enumerate(prior):
        cx = x + i * (cw + 0.12)
        rect(s, cx, y, cw, 0.84, fill=BG, line=LINE)
        rect(s, cx, y, cw, 0.045, fill=col)
        txt(s, cx + 0.14, y + 0.12, cw - 0.28, 0.24, titulo, size=10, bold=True, color=col)
        txt(s, cx + 0.14, y + 0.38, cw - 0.28, 0.42, det, size=9, color=INK)

    # 4 · metas
    y = 3.46
    x, w = _banda(s, y, 1.94, "METAS\n2028")
    rect(s, x, y, w, 1.94, fill=WHITE, line=LINE)
    cw = (w - 0.30) / 2
    for i, (grupo, metas, col) in enumerate([
            ("Velocidad de crecimiento · ampliar horizontes", METAS_CRECIMIENTO, STEEL),
            ("Velocidad de consolidación del CORE", METAS_CORE, NAVY)]):
        cx = x + 0.10 + i * (cw + 0.10)
        rect(s, cx, y + 0.08, cw, 0.26, fill=col)
        txt(s, cx + 0.10, y + 0.12, cw - 0.20, 0.20, grupo.upper(), size=8, bold=True, color=WHITE)
        for j, (tema, meta, det) in enumerate(metas):
            my = y + 0.40 + j * 0.375
            n = i * 4 + j + 1
            o = oval(s, cx + 0.04, my + 0.03, 0.22, fill=col)
            shape_text(o, str(n), size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER, pad=0)
            txt_inline(s, cx + 0.32, my - 0.02, cw - 0.40, 0.20,
                       [(f"{tema}:  ", {"size": 8.5, "bold": True, "color": GREY}),
                        (meta, {"size": 9.5, "bold": True, "color": INK})])
            txt(s, cx + 0.32, my + 0.16, cw - 0.40, 0.18, det, size=7.5, color=GREY)

    # 5 · proceso de generación de valor
    y = 5.52
    x, w = _banda(s, y, 0.60, "PROCESO DE\nGENERACIÓN\nDE VALOR")
    cw = w / 4
    for i, (titulo, det) in enumerate(CADENA_VALOR):
        cx = x + i * cw
        shape = MSO_SHAPE.PENTAGON if i == 0 else MSO_SHAPE.CHEVRON
        sh = rect(s, cx, y, cw + 0.10, 0.60, fill=(NAVY if i in (0, 3) else STEEL), shape=shape)
        shape_text(sh, [(titulo, {"size": 9.5, "bold": True, "color": WHITE}),
                        (det, {"size": 7, "italic": True, "color": ICE})],
                   align=PP_ALIGN.CENTER, pad=0.14)

    # 6 · foco y habilitadores
    y = 6.24
    x, w = _banda(s, y, 0.82, "FOCO Y\nHABILITADORES")
    cw = (w - 0.36) / 4
    for i, (titulo, det) in enumerate(HABILITADORES):
        cx = x + i * (cw + 0.12)
        rect(s, cx, y, cw, 0.52, fill=BG, line=LINE)
        txt(s, cx + 0.10, y + 0.03, cw - 0.20, 0.18, titulo, size=8, bold=True, color=INK)
        txt(s, cx + 0.10, y + 0.21, cw - 0.20, 0.30, det, size=6.5, color=GREY)
    rect(s, x, y + 0.56, w, 0.26, fill=STEEL)
    txt(s, x + 0.20, y + 0.59, w - 0.40, 0.20,
        "Sistema de Gestión Estratégica enfocado en el equilibrio entre fortalecimiento del "
        "CORE y crecimiento", size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    source_note(s, "Propósito, misión, motores, cadena de valor y habilitadores: conservados del "
                   "deck de Gestión Estratégica REDCO, junio 2024 (diapositiva 38). Las ocho metas "
                   "son nuevas y su línea base sale de pilar_a. La meta 2024 («sostener ~12 MUS$ y "
                   "~18 MUS$ al 2026») no se cumplió: 2025 cerró en 9,5 MUSD emitidos.", y=7.14)


def s22_cierre(prs):
    cover(blank(prs),
          "Gestión Estratégica REDCO",
          "Ocho metas al 2028, con línea base medida y trazable",
          "Capítulos 2 – 4  ·  Actualización 2026",
          "REDCO · Mining Consultants")


# ========================================================================= main

def main():
    d = cargar()
    prs = nueva_presentacion()

    s01_portada(prs)
    s02_agenda(prs)
    s03_divisor_cap2(prs)
    s04_contexto(prs, d)
    s05_crecimiento(prs, d)
    s06_ventas(prs, d)
    s07_rentabilidad_pais(prs, d)
    s08_rentabilidad_nivel(prs, d)
    s09_mercado_etc(prs)
    s10a_cobre(prs)
    s10b_cobre(prs)
    s11_fundamentos(prs)
    s12_core(prs)
    s13_divisor_cap3(prs)
    s14_elementos(prs)
    s15_flujo_decisiones(prs)
    s16_distancia_core(prs)
    s17_usa(prs, d)
    s18_probabilidades(prs)
    s19_divisor_cap4(prs)
    s20_motores(prs)
    s21_proposito_metas(prs)
    s22_cierre(prs)

    prs.save(OUT)
    print(f"OK  {OUT}  ·  {len(prs.slides.__iter__.__self__._sldIdLst)} slides")


if __name__ == "__main__":
    main()
