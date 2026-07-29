"""Vocabulario visual REDCO para presentaciones .pptx.

Implementa el `Estandar_Presentaciones_PPT_REDCO(1).md` replicando la construcción
ya validada en `REDCO_Presentacion_Historico_rev1.pptx`: paleta exacta, escala
tipográfica, encabezado con viñeta circular, caja de insight y pie de fuente.

Tipografía: el estándar (§4.1) pide declarar Lexend Deca en el tema con fallback
explícito a una fuente segura. Lexend Deca no está instalada en el entorno, así
que el tema la declara y los runs usan Calibri — igual que rev1 — para que el
render sea fiel y la revisión visual sea válida.
"""

from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------- paleta (§5)
NAVY = RGBColor(0x1E, 0x27, 0x61)   # azul medianoche · dominante
STEEL = RGBColor(0x2C, 0x4A, 0x7C)  # azul acero · secundario
ICE = RGBColor(0xCA, 0xDC, 0xFC)    # celeste hielo · acento claro
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x23, 0x27, 0x33)    # texto cuerpo
GREY = RGBColor(0x6B, 0x72, 0x80)   # metadatos, labels, notas
LINE = RGBColor(0xE2, 0xE5, 0xEA)   # bordes
BG = RGBColor(0xF4, 0xF6, 0xF9)     # fondo de tarjeta

# estados — uso exclusivo en leyendas MECE de 3 estados (§5)
GREEN = RGBColor(0x2E, 0x7D, 0x5B)
AMBER = RGBColor(0xC8, 0x86, 0x2D)
RED = RGBColor(0xC8, 0x10, 0x2E)

# Rampa categórica de marca. §5 reserva verde/ámbar/rojo para leyendas de estado,
# así que las series de un gráfico NUNCA los usan: se diferencian dentro del azul.
MID = RGBColor(0x3E, 0x6F, 0xB0)
SKY = RGBColor(0x7F, 0xA3, 0xD4)
SLATE = RGBColor(0xA8, 0xAE, 0xB8)
SERIES_RAMP = [NAVY, STEEL, MID, SKY, ICE, GREY, SLATE, LINE]

FONT = "Calibri"
BRAND_FONT = "Lexend Deca"

# ------------------------------------------------------------- geometría (§3)
SLIDE_W, SLIDE_H = 13.333, 7.5
MARGIN = 0.5          # margen de seguridad en los 4 bordes
GAP = 0.3             # separación mínima entre bloques
CONTENT_TOP = 1.62    # primera línea útil bajo el encabezado
FOOTER_Y = 7.03       # línea del pie de fuente


def _srgb(hexstr):
    return RGBColor.from_string(hexstr.lstrip("#"))


# --------------------------------------------------------------- primitivas

def txt(slide, x, y, w, h, text, size=12, bold=False, color=INK, align=PP_ALIGN.LEFT,
        italic=False, anchor=MSO_ANCHOR.TOP, spacing=None, wrap=True, font=FONT):
    """Caja de texto. `text` puede ser str o lista de (texto, kwargs) por párrafo."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0

    paras = text if isinstance(text, list) else [(text, {})]
    for i, item in enumerate(paras):
        line, opts = item if isinstance(item, tuple) else (item, {})
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get("align", align)
        sp = opts.get("spacing", spacing)
        if sp is not None:
            p.space_after = Pt(sp)
        run = p.add_run()
        run.text = line
        f = run.font
        f.name = opts.get("font", font)
        f.size = Pt(opts.get("size", size))
        f.bold = opts.get("bold", bold)
        f.italic = opts.get("italic", italic)
        f.color.rgb = opts.get("color", color)
    return box


def txt_inline(slide, x, y, w, h, runs, size=11, color=INK, align=PP_ALIGN.LEFT,
               anchor=MSO_ANCHOR.TOP):
    """Una sola línea con varios runs de formato distinto.

    `txt()` crea un párrafo por elemento; esto los mantiene en el MISMO párrafo,
    que es lo que hace falta para pares etiqueta+valor.
    """
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    for text, opts in runs:
        run = p.add_run()
        run.text = text
        f = run.font
        f.name = opts.get("font", FONT)
        f.size = Pt(opts.get("size", size))
        f.bold = opts.get("bold", False)
        f.italic = opts.get("italic", False)
        f.color.rgb = opts.get("color", color)
    return box


def rect(slide, x, y, w, h, fill=None, line=None, line_w=0.75, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    s.shadow.inherit = False
    if s.has_text_frame:
        s.text_frame.word_wrap = True
    return s


def oval(slide, x, y, d, fill=NAVY, line=None):
    return rect(slide, x, y, d, d, fill=fill, line=line, shape=MSO_SHAPE.OVAL)


def line_h(slide, x, y, w, color=LINE, width=0.75):
    s = slide.shapes.add_connector(1, Inches(x), Inches(y), Inches(x + w), Inches(y))
    s.line.color.rgb = color
    s.line.width = Pt(width)
    return s


def shape_text(shape, text, size=12, bold=False, color=INK, align=PP_ALIGN.LEFT,
               italic=False, anchor=MSO_ANCHOR.MIDDLE, pad=0.06, spacing=0):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(pad)
    tf.margin_top = tf.margin_bottom = Inches(pad / 2)
    paras = text if isinstance(text, list) else [(text, {})]
    for i, item in enumerate(paras):
        line, opts = item if isinstance(item, tuple) else (item, {})
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = opts.get("align", align)
        p.space_after = Pt(opts.get("spacing", spacing))
        run = p.add_run()
        run.text = line
        f = run.font
        f.name = FONT
        f.size = Pt(opts.get("size", size))
        f.bold = opts.get("bold", bold)
        f.italic = opts.get("italic", italic)
        f.color.rgb = opts.get("color", color)
    return shape


# ------------------------------------------------------- componentes (§7, §8)

def header(slide, title, synthesis=None, brand=True):
    """Encabezado de slide de contenido: viñeta circular + título + síntesis (§7.1-7.2).

    El tamaño del título se ajusta al largo para que SIEMPRE quepa en una línea:
    si se desborda a dos, pisa la frase de síntesis que va justo debajo.
    """
    if brand:
        txt(slide, SLIDE_W - 3.2 - MARGIN, 0.30, 3.2, 0.22, "REDCO  ·  Mining Consultants",
            size=9, color=GREY, align=PP_ALIGN.RIGHT)
    # viñeta circular azul acero con punto blanco interior — motivo de marca
    oval(slide, MARGIN, 0.62, 0.20, fill=STEEL)
    oval(slide, MARGIN + 0.068, 0.688, 0.064, fill=WHITE)

    n = len(title)
    size = 26 if n <= 62 else (23 if n <= 74 else (20 if n <= 88 else 18))
    txt(slide, MARGIN + 0.34, 0.53 + (26 - size) * 0.006,
        SLIDE_W - 2 * MARGIN - 0.34, 0.50, title, size=size, bold=True, color=INK)
    if synthesis:
        txt(slide, MARGIN + 0.34, 1.06, SLIDE_W - 2 * MARGIN - 0.34, 0.40,
            synthesis, size=13, italic=True, color=GREY)
    # espacio reservado para el logo (§6) — esquina superior izquierda
    return slide


def source_note(slide, text, y=FOOTER_Y):
    txt(slide, MARGIN, y, SLIDE_W - 2 * MARGIN, 0.34, text, size=8, color=GREY)


def insight(slide, x, y, w, h, label, headline, body, accent=NAVY):
    """Caja de insight: label en mayúscula + cifra grande + texto (patrón de rev1)."""
    rect(slide, x, y, w, h, fill=BG, line=LINE)
    rect(slide, x, y, 0.05, h, fill=accent)
    txt(slide, x + 0.20, y + 0.13, w - 0.40, 0.20, label, size=8.5, bold=True, color=GREY)
    cur = y + 0.34
    if headline:
        txt(slide, x + 0.20, cur, w - 0.40, 0.42, headline, size=25, bold=True, color=accent)
        cur += 0.46
    # si el alto disponible se queda corto, se baja el cuerpo antes que desbordar la caja
    disponible = h - (cur - y) - 0.10
    size = 10.5 if disponible >= 0.60 else (9.5 if disponible >= 0.42 else 8.5)
    txt(slide, x + 0.20, cur, w - 0.40, max(disponible, 0.20), body, size=size, color=INK)


def card(slide, x, y, w, h, number, title, bullets, owner=None):
    """Tarjeta numerada (§8) — numeración SIEMPRE azul acero, nunca multicolor."""
    rect(slide, x, y, w, h, fill=BG, line=LINE)
    oval(slide, x + 0.18, y + 0.18, 0.30, fill=STEEL)
    shape_text(slide.shapes[-1], str(number), size=11.5, bold=True, color=WHITE,
               align=PP_ALIGN.CENTER, pad=0)
    txt(slide, x + 0.58, y + 0.17, w - 0.76, 0.46, title, size=13, bold=True, color=INK)
    body_y = y + 0.17 + (0.30 if len(title) < 34 else 0.52)
    items = [(f"•  {b}", {"size": 11, "spacing": 3}) for b in bullets]
    txt(slide, x + 0.20, body_y, w - 0.40, h - (body_y - y) - 0.14, items, color=INK)
    if owner:
        txt(slide, x + 0.20, y + h - 0.34, w - 0.40, 0.24,
            [("RESPONSABLE  ", {"size": 8, "bold": True, "color": GREY}),
             (owner, {"size": 10, "bold": True, "color": INK})])


def status_legend(slide, y=None, x=None):
    """Leyenda de estado: exactamente 3 estados, siempre al pie (§8)."""
    y = FOOTER_Y - 0.34 if y is None else y
    x = MARGIN if x is None else x
    for i, (color, label) in enumerate(
            [(GREEN, "En línea con lo esperado"),
             (AMBER, "Requiere atención puntual"),
             (RED, "Crítico / requiere intervención")]):
        cx = x + i * 3.15
        oval(slide, cx, y + 0.045, 0.11, fill=color)
        txt(slide, cx + 0.19, y, 2.85, 0.22, label, size=9, color=GREY)


def cover(slide, title, subtitle, meta, footer):
    """Portada / cierre: fondo azul medianoche a sangre completa (§7)."""
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)
    txt(slide, 0, 1.30, SLIDE_W, 0.34, "REDCO  ·  MINING CONSULTANTS",
        size=12, bold=True, color=ICE, align=PP_ALIGN.CENTER)
    txt(slide, 1.4, 2.30, SLIDE_W - 2.8, 1.70, title,
        size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(slide, SLIDE_W / 2 - 0.85, 4.18, 1.70, 0.030, fill=ICE)
    txt(slide, 1.4, 4.52, SLIDE_W - 2.8, 0.40, subtitle,
        size=15, color=ICE, align=PP_ALIGN.CENTER)
    txt(slide, 1.4, 5.06, SLIDE_W - 2.8, 0.30, meta,
        size=11, color=ICE, align=PP_ALIGN.CENTER)
    txt(slide, 1.4, SLIDE_H - 1.00, SLIDE_W - 2.8, 0.28, footer,
        size=9, color=ICE, align=PP_ALIGN.CENTER)


def divider(slide, number, title, synthesis):
    """Divisor de capítulo."""
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)
    rect(slide, 0, 0, 0.26, SLIDE_H, fill=STEEL)
    txt(slide, 1.5, 2.44, 3.0, 1.40, number, size=88, bold=True, color=ICE)
    # mismo criterio que header(): el título no puede pasar a dos líneas o pisa
    # la línea divisoria y la frase de síntesis que van debajo
    t_size = 40 if len(title) <= 30 else (34 if len(title) <= 42 else 28)
    txt(slide, 1.5, 3.72 + (40 - t_size) * 0.008, SLIDE_W - 3.0, 0.80,
        title, size=t_size, bold=True, color=WHITE)
    rect(slide, 1.5, 4.66, 1.40, 0.030, fill=ICE)
    txt(slide, 1.5, 4.94, SLIDE_W - 3.4, 0.70, synthesis, size=14, color=ICE)
    txt(slide, SLIDE_W - 3.2 - MARGIN, SLIDE_H - 0.72, 3.2, 0.24,
        "REDCO  ·  Mining Consultants", size=9, color=ICE, align=PP_ALIGN.RIGHT)


# ------------------------------------------------------------------ gráficos

def _style_chart(chart, font_size=9, legend=False, legend_pos=XL_LEGEND_POSITION.BOTTOM):
    chart.font.size = Pt(font_size)
    chart.font.name = FONT
    chart.font.color.rgb = GREY
    # con una sola serie, el renderer inventa un título con el nombre de la serie
    chart.has_title = False
    chart.has_legend = legend
    if legend:
        chart.legend.position = legend_pos
        chart.legend.include_in_layout = False
    return chart


def add_chart(slide, kind, x, y, w, h, categories, series, colors=None,
              legend=False, gap=60, overlap=None, font_size=9, point_colors=None,
              y_min=None, y_max=None):
    cd = CategoryChartData()
    cd.categories = categories
    for name, vals in series:
        cd.add_series(name, vals)
    gf = slide.shapes.add_chart(kind, Inches(x), Inches(y), Inches(w), Inches(h), cd)
    chart = gf.chart
    _style_chart(chart, font_size=font_size, legend=legend)
    plot = chart.plots[0]
    try:
        plot.gap_width = gap
        if overlap is not None:
            plot.overlap = overlap
    except (AttributeError, ValueError):
        pass
    if colors:
        for s, col in zip(chart.series, colors):
            s.format.fill.solid()
            s.format.fill.fore_color.rgb = col
            s.format.line.fill.background()
    if point_colors:
        # dona/torta: el color va por PUNTO, no por serie
        pts = chart.plots[0].series[0].points
        for pt, col in zip(pts, point_colors):
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = col
            pt.format.line.color.rgb = WHITE
            pt.format.line.width = Pt(1)
    if y_min is not None or y_max is not None:
        try:
            if y_min is not None:
                chart.value_axis.minimum_scale = y_min
            if y_max is not None:
                chart.value_axis.maximum_scale = y_max
        except (AttributeError, ValueError):
            pass
    return chart
